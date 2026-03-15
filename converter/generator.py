import json
import os
import shutil
import tempfile
from dataclasses import replace
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .adapters import MinerUAdapter, MinerUPageData, OCRAdapter
from .ir import ElementIR, ImageIR, TextIR, TextRunIR, build_page_ir, materialize_text_runs_for_elements, validate_ir_elements
from .ir_merge import merge_ir_elements
from .ocr_merge import PaddleOCREngine
from .utils import extract_background_color, extract_font_color, fill_bbox_with_bg


class PageContext:
    def __init__(self, page_index, page_image, coords, slide, temp_dir: Path | None = None, debug_dir: Path | None = None):
        self.page_index = int(page_index)
        self.slide = slide
        self.original_image = page_image.copy()
        self.background_image = page_image.copy()
        self.coords = coords
        self.elements = []
        self.stage_page_irs = {}
        self.temp_dir = Path(temp_dir) if temp_dir else Path(tempfile.gettempdir())
        self.debug_dir = Path(debug_dir) if debug_dir else None

    def add_element_bbox_for_cleanup(self, bbox, margin_px=0):
        """Register a bounding box to be inpainted on the background image."""
        if bbox:
            px_box = [int(v * (self.coords['img_w'] / self.coords['json_w'] if i % 2 == 0 else self.coords['img_h'] / self.coords['json_h'])) for i, v in enumerate(bbox)]
            if margin_px > 0:
                x1, y1, x2, y2 = px_box
                px_box = [x1 - margin_px, y1 - margin_px, x2 + margin_px, y2 + margin_px]

            fill_bbox_with_bg(self.background_image, px_box)

    def add_processed_element(self, elem_type, data):
        """Store a fully processed element ready for rendering."""
        self.elements.append({'type': elem_type, 'data': data})

    def register_stage_page_ir(self, stage, page_ir):
        self.stage_page_irs[str(stage)] = page_ir

    def _extract_stage_text_bboxes(self, stage):
        page_ir = self.stage_page_irs.get(stage)
        if page_ir is None:
            return []
        return [
            elem.bbox
            for elem in page_ir.elements
            if isinstance(elem, TextIR) and elem.bbox
        ]

    def generate_debug_images(self, generator_instance):
        """Generate and save debug images for the page."""
        if self.debug_dir is None:
            return

        page_index = self.page_index
        self.debug_dir.mkdir(parents=True, exist_ok=True)

        cv2.imwrite(str(self.debug_dir / f"page_{page_index}_original.png"), cv2.cvtColor(self.original_image, cv2.COLOR_RGB2BGR))

        stage_outputs = [
            ("mineru_original", self.debug_dir / f"page_{page_index}_mineru_original_boxes.png"),
            ("ocr_before_refined_elements", self.debug_dir / f"page_{page_index}_ocr_before_refined_elements.png"),
            ("ocr_after_refined_elements", self.debug_dir / f"page_{page_index}_ocr_after_refined_elements.png"),
            ("merged_final", self.debug_dir / f"page_{page_index}_merged_final_boxes.png"),
        ]

        for stage_name, output_path in stage_outputs:
            stage_bboxes = self._extract_stage_text_bboxes(stage_name)
            generator_instance._draw_text_bboxes_for_page(
                self.original_image,
                stage_bboxes,
                self.coords,
                str(output_path),
            )

        text_bboxes = [
            elem['data']['bbox']
            for elem in self.elements
            if elem.get('type') == 'text' and elem.get('data', {}).get('bbox')
        ]
        generator_instance._draw_text_bboxes_for_page(
            self.original_image,
            text_bboxes,
            self.coords,
            str(self.debug_dir / f"page_{page_index}_text_boxes.png"),
        )

    def render_to_slide(self, generator_instance):
        """Render all processed elements onto the PowerPoint slide."""
        # 1. Render the cleaned background
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        fd, bg_path = tempfile.mkstemp(prefix="temp_bg_", suffix=".png", dir=str(self.temp_dir))
        os.close(fd)
        try:
            cv2.imwrite(bg_path, cv2.cvtColor(self.background_image, cv2.COLOR_RGB2BGR))
            w_pts, h_pts = generator_instance.prs.slide_width, generator_instance.prs.slide_height
            self.slide.shapes.add_picture(bg_path, Pt(0), Pt(0), w_pts, h_pts)
        finally:
            if os.path.exists(bg_path):
                os.remove(bg_path)

        # 2. Render all image elements first
        for elem in self.elements:
            if elem['type'] == 'image':
                generator_instance._add_picture_from_bbox(
                    self.slide,
                    elem['data']['bbox'],
                    self.original_image,
                    self.coords,
                    elem['data'].get('image_crop'),
                )

        # 3. Render all text elements on top
        for elem in self.elements:
            if elem['type'] == 'text':
                generator_instance._render_text_from_data(self.slide, elem['data'])


TEXT_CLEANUP_MARGIN_RATIO = 0.05
TEXT_CLEANUP_MIN_MARGIN_PX = 1


class PPTGenerator:
    @staticmethod
    def _safe_float(value):
        if value is None:
            return None
        try:
            parsed = float(value)
        except (TypeError, ValueError):
            return None
        return parsed if parsed > 0 else None

    @classmethod
    def _effective_font_size_from_run(cls, elem: TextIR, run: TextRunIR):
        run_style = run.style or {}
        elem_style = elem.style or {}
        style_font_size = cls._safe_float(run_style.get("font_size"))
        if style_font_size is None:
            style_font_size = cls._safe_float(elem_style.get("font_size"))
        if style_font_size is not None:
            return style_font_size

        run_bbox = run.bbox or elem.bbox
        if isinstance(run_bbox, (list, tuple)) and len(run_bbox) == 4:
            try:
                height = float(run_bbox[3]) - float(run_bbox[1])
            except (TypeError, ValueError):
                return None
            if height > 0:
                return height
        return None

    @classmethod
    def _effective_font_size_from_element(cls, elem: TextIR):
        elem_style = elem.style or {}
        style_font_size = cls._safe_float(elem_style.get("font_size"))
        if style_font_size is not None:
            return style_font_size

        bbox = elem.bbox
        if isinstance(bbox, (list, tuple)) and len(bbox) == 4:
            try:
                height = float(bbox[3]) - float(bbox[1])
            except (TypeError, ValueError):
                return None
            if height > 0:
                return height
        return None

    @staticmethod
    def _ratio_distance_to_center(size_value: float, center_value: float) -> float:
        if size_value <= 0 or center_value <= 0:
            return float("inf")
        return max(size_value, center_value) / min(size_value, center_value)

    def _assign_groups_by_center_distance(self, samples: list[tuple[float, int]], ratio_threshold: float) -> list[list[tuple[float, int]]]:
        if not samples:
            return []

        ordered_samples = sorted(samples, key=lambda item: (item[0], item[1]))
        clusters: list[list[tuple[float, int]]] = []

        for sample in ordered_samples:
            candidate_clusters: list[tuple[float, int]] = []
            for cluster_index, cluster in enumerate(clusters):
                candidate = cluster + [sample]
                center_value = float(np.median([entry[0] for entry in candidate]))
                max_distance = max(
                    self._ratio_distance_to_center(entry[0], center_value)
                    for entry in candidate
                )
                if max_distance <= ratio_threshold:
                    candidate_clusters.append((max_distance, cluster_index))

            if candidate_clusters:
                _, best_cluster_index = min(candidate_clusters, key=lambda item: (item[0], item[1]))
                clusters[best_cluster_index].append(sample)
            else:
                clusters.append([sample])

        return [sorted(cluster, key=lambda item: (item[0], item[1])) for cluster in clusters]

    def _groups_within_center_threshold(self, groups: list[list[tuple[float, int]]], ratio_threshold: float) -> bool:
        for group in groups:
            if not group:
                continue
            center_value = float(np.median([entry[0] for entry in group]))
            for size_value, _elem_index in group:
                if self._ratio_distance_to_center(size_value, center_value) > ratio_threshold:
                    return False
        return True

    def _optimize_groups_with_kmeans(
        self,
        samples: list[tuple[float, int]],
        seed_groups: list[list[tuple[float, int]]],
        ratio_threshold: float,
    ) -> list[list[tuple[float, int]]]:
        if not seed_groups:
            return []

        ordered_samples = sorted(samples, key=lambda item: (item[0], item[1]))
        sample_count = len(ordered_samples)
        k = len(seed_groups)

        if sample_count <= 1 or k <= 1:
            return seed_groups
        if k >= sample_count:
            return [[sample] for sample in ordered_samples]

        log_values = np.array([np.log(sample[0]) for sample in ordered_samples], dtype=float)
        seed_centers = np.array(
            [
                float(np.median(np.log([entry[0] for entry in group])))
                for group in seed_groups
                if group
            ],
            dtype=float,
        )
        if seed_centers.size != k:
            return seed_groups

        centers = np.sort(seed_centers)

        for _ in range(30):
            distances = np.abs(log_values[:, None] - centers[None, :])
            labels = np.argmin(distances, axis=1)

            new_centers = centers.copy()
            for cluster_index in range(k):
                cluster_mask = labels == cluster_index
                if np.any(cluster_mask):
                    new_centers[cluster_index] = float(np.mean(log_values[cluster_mask]))
                else:
                    farthest_idx = int(np.argmax(np.min(distances, axis=1)))
                    new_centers[cluster_index] = log_values[farthest_idx]

            if np.allclose(new_centers, centers, atol=1e-7):
                centers = new_centers
                break
            centers = new_centers

        final_distances = np.abs(log_values[:, None] - centers[None, :])
        final_labels = np.argmin(final_distances, axis=1)

        optimized_groups: list[list[tuple[float, int]]] = []
        for cluster_index in range(k):
            cluster = [
                ordered_samples[sample_index]
                for sample_index in range(sample_count)
                if final_labels[sample_index] == cluster_index
            ]
            if cluster:
                optimized_groups.append(sorted(cluster, key=lambda item: (item[0], item[1])))

        optimized_groups = sorted(optimized_groups, key=lambda cluster: (cluster[0][0], cluster[0][1]))
        if self._groups_within_center_threshold(optimized_groups, ratio_threshold):
            return optimized_groups
        return seed_groups

    @staticmethod
    def _median_rounded_int(values: list[float]) -> int | None:
        if not values:
            return None
        return int(round(float(np.median(values))))

    @classmethod
    def _effective_bold_from_text_element(cls, elem: TextIR) -> bool:
        runs = elem.text_runs if isinstance(elem.text_runs, list) else []
        elem_style = elem.style or {}
        if not runs:
            return bool(elem_style.get("bold", False))

        run_bolds = [
            bool((run.style or {}).get("bold", elem_style.get("bold", False)))
            for run in runs
        ]
        true_count = sum(1 for flag in run_bolds if flag)
        false_count = len(run_bolds) - true_count
        if true_count == false_count:
            return bool(elem_style.get("bold", False))
        return true_count > false_count

    @classmethod
    def _effective_font_size_from_text_element(cls, elem: TextIR):
        runs = elem.text_runs if isinstance(elem.text_runs, list) else []
        run_sizes = [
            cls._effective_font_size_from_run(elem, run)
            for run in runs
        ]
        run_sizes = [size for size in run_sizes if size is not None]
        if run_sizes:
            return float(np.median(run_sizes))
        return cls._effective_font_size_from_element(elem)

    @classmethod
    def _normalize_text_element_internal_format(cls, elem: TextIR) -> TextIR | None:
        runs = elem.text_runs if isinstance(elem.text_runs, list) else []
        elem_style = elem.style or {}

        if not runs:
            unified_bold = bool(elem_style.get("bold", False))
            unified_size = cls._effective_font_size_from_element(elem)
            if unified_size is None:
                return None

            new_elem_style = dict(elem_style)
            new_elem_style["bold"] = unified_bold
            new_elem_style["font_size"] = float(unified_size)
            return replace(elem, style=new_elem_style)

        run_bolds: list[bool] = []
        run_sizes: list[float] = []
        for run in runs:
            run_style = run.style or {}
            run_bolds.append(bool(run_style.get("bold", elem_style.get("bold", False))))
            size_value = cls._effective_font_size_from_run(elem, run)
            if size_value is None:
                return None
            run_sizes.append(size_value)

        if not run_bolds or not run_sizes:
            return None

        if any(flag != run_bolds[0] for flag in run_bolds):
            return None

        min_size = min(run_sizes)
        max_size = max(run_sizes)
        if min_size <= 0 or (max_size / min_size) > 1.3:
            return None

        unified_bold = run_bolds[0]
        unified_size = float(np.median(run_sizes))

        new_elem_style = dict(elem_style)
        new_elem_style["bold"] = unified_bold
        new_elem_style["font_size"] = float(unified_size)

        new_runs: list[TextRunIR] = []
        for run in runs:
            run_style = dict(run.style or {})
            run_style["bold"] = unified_bold
            run_style["font_size"] = float(unified_size)
            new_runs.append(replace(run, style=run_style))

        return replace(elem, style=new_elem_style, text_runs=new_runs)

    def _apply_font_size_to_text_element(self, elem: TextIR, target_size: int) -> TextIR:
        target_value = float(target_size)

        elem_style = dict(elem.style or {})
        elem_style["font_size"] = target_value
        updated = replace(elem, style=elem_style)

        runs = updated.text_runs if isinstance(updated.text_runs, list) else []
        if not runs:
            return updated

        new_runs: list[TextRunIR] = []
        for run in runs:
            run_style = dict(run.style or {})
            run_style["font_size"] = target_value
            new_runs.append(replace(run, style=run_style))

        return replace(updated, text_runs=new_runs)

    def _normalize_page_text_font_sizes(self, elements: list[ElementIR]) -> list[ElementIR]:
        text_elements = [elem for elem in elements if isinstance(elem, TextIR)]
        if not text_elements:
            return elements

        normalized_pairs: list[tuple[TextIR, bool]] = []
        for elem in text_elements:
            normalized = self._normalize_text_element_internal_format(elem)
            if normalized is None:
                normalized_pairs.append((elem, False))
            else:
                normalized_pairs.append((normalized, True))

        bucket_samples: dict[bool, list[tuple[float, int]]] = {
            False: [],
            True: [],
        }
        for elem_index, (elem, is_eligible) in enumerate(normalized_pairs):
            if not is_eligible:
                continue
            size_value = self._effective_font_size_from_text_element(elem)
            if size_value is None:
                continue
            bold_flag = self._effective_bold_from_text_element(elem)
            bucket_samples[bold_flag].append((size_value, elem_index))

        updated_text_elements = [elem for elem, _eligible in normalized_pairs]
        ratio_threshold = 1.3

        for _bold_flag, samples in bucket_samples.items():
            if not samples:
                continue

            grouped_samples = self._assign_groups_by_center_distance(samples, ratio_threshold)
            grouped_samples = self._optimize_groups_with_kmeans(samples, grouped_samples, ratio_threshold)

            for group in grouped_samples:
                if len(group) < 2:
                    continue

                median_size = self._median_rounded_int([entry[0] for entry in group])
                if median_size is None:
                    continue

                for _size, elem_index in group:
                    updated_text_elements[elem_index] = self._apply_font_size_to_text_element(
                        updated_text_elements[elem_index],
                        median_size,
                    )

        text_iter = iter(updated_text_elements)
        rebuilt_elements: list[ElementIR] = []
        for elem in elements:
            if isinstance(elem, TextIR):
                rebuilt_elements.append(next(text_iter))
            else:
                rebuilt_elements.append(elem)
        return rebuilt_elements
    def __init__(
        self,
        output_path,
        remove_watermark=True,
        ocr_engine=None,
        ocr_device_policy="auto",
        ocr_model_root=None,
        ocr_model_variant="auto",
        ocr_offline_only=False,
        text_cleanup_margin_ratio=None,
        temp_dir=None,
        debug_dir=None,
    ):
        self.prs = Presentation()
        self.output_path = output_path
        self.temp_dir = Path(temp_dir) if temp_dir else Path(tempfile.gettempdir())
        self.debug_dir = Path(debug_dir) if debug_dir else None
        self.remove_watermark = remove_watermark
        self.ocr_engine = ocr_engine
        self.ocr_device_policy = ocr_device_policy
        self.ocr_model_root = ocr_model_root
        self.ocr_model_variant = ocr_model_variant
        self.ocr_offline_only = ocr_offline_only
        if text_cleanup_margin_ratio is None:
            text_cleanup_margin_ratio = TEXT_CLEANUP_MARGIN_RATIO
        self.text_cleanup_margin_ratio = float(text_cleanup_margin_ratio)
        self.debug_images = False # Will be set in process_page
        for i in range(len(self.prs.slides) - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]

    def cap_size(self, w_pts, h_pts):
        MAX_PTS = 56 * 72
        if w_pts > MAX_PTS or h_pts > MAX_PTS:
            scale = MAX_PTS / max(w_pts, h_pts)
            w_pts, h_pts = w_pts * scale, h_pts * scale
        return w_pts, h_pts

    def set_slide_size(self, width_px, height_px, dpi=72):
        w_pts, h_pts = self.cap_size(width_px * 72 / dpi, height_px * 72 / dpi)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)

    def add_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _get_bbox_intersection(self, bbox1, bbox2):
        x1, y1 = max(bbox1[0], bbox2[0]), max(bbox1[1], bbox2[1])
        x2, y2 = min(bbox1[2], bbox2[2]), min(bbox1[3], bbox2[3])
        return [x1, y1, x2, y2] if x1 < x2 and y1 < y2 else None

    def _has_bbox_overlap(self, bbox1, bbox2):
        return self._get_bbox_intersection(bbox1, bbox2) is not None

    def _to_px_bbox(self, bbox, coords):
        return [
            int(bbox[0] * coords['img_w'] / coords['json_w']),
            int(bbox[1] * coords['img_h'] / coords['json_h']),
            int(bbox[2] * coords['img_w'] / coords['json_w']),
            int(bbox[3] * coords['img_h'] / coords['json_h']),
        ]

    def _create_textbox(self, slide, bbox, coords):
        x1, y1, x2, y2 = bbox
        return slide.shapes.add_textbox(
            Pt(x1 * coords['scale_x']), Pt(y1 * coords['scale_y']),
            Pt((x2 - x1) * coords['scale_x']), Pt((y2 - y1) * coords['scale_y'])
        )

    def _draw_text_bboxes_for_page(self, image, text_bboxes, coords, output_path):
        """Draw text-level bounding boxes for debugging."""
        debug_img = image.copy()
        for bbox in text_bboxes:
            px_box = [
                int(bbox[0] * coords['img_w'] / coords['json_w']),
                int(bbox[1] * coords['img_h'] / coords['json_h']),
                int(bbox[2] * coords['img_w'] / coords['json_w']),
                int(bbox[3] * coords['img_h'] / coords['json_h'])
            ]
            cv2.rectangle(debug_img, (px_box[0], px_box[1]), (px_box[2], px_box[3]), (0, 255, 0), 2)
        cv2.imwrite(output_path, cv2.cvtColor(debug_img, cv2.COLOR_RGB2BGR))

    def _build_text_runs_from_ir_runs(self, context, elem: TextIR, ir_runs: list[TextRunIR]):
        bbox = elem.bbox
        if not bbox or not ir_runs:
            return [], False

        ordered_runs = sorted(
            ir_runs,
            key=lambda run: (
                int(run.line_index),
                (run.bbox or bbox)[0],
                (run.bbox or bbox)[1],
            ),
        )

        runs = []
        current_line_index = None
        line_indexes = set()

        for run in ordered_runs:
            run_text = str(run.text).replace("\\%", "%")
            if not run_text:
                continue

            run_bbox = run.bbox or bbox
            line_index = int(run.line_index)
            line_indexes.add(line_index)

            if current_line_index is not None and line_index != current_line_index and runs:
                prev_font = runs[-1].get("font", {})
                runs.append({"text": "\n", "font": prev_font})

            run_px_bbox = self._to_px_bbox(run_bbox, context.coords)
            bg_color = extract_background_color(context.original_image, run_px_bbox)
            color, _, _ = extract_font_color(context.original_image, run_px_bbox, bg_color)

            run_style = run.style or {}
            elem_style = elem.style or {}
            style_font_size = run_style.get("font_size") or elem_style.get("font_size")
            if style_font_size:
                font_size_pts = max(6.0, float(style_font_size) * context.coords['scale_y'])
            else:
                font_size_pts = max(6.0, (run_bbox[3] - run_bbox[1]) * context.coords['scale_y'])

            font_info = {
                "name": "Microsoft YaHei",
                "size": Pt(int(font_size_pts)),
                "bold": bool(run_style.get("bold", elem_style.get("bold", False))),
                "color": RGBColor(*color),
            }

            runs.append({"text": run_text, "font": font_info})
            current_line_index = line_index

        if not runs:
            return [], False

        is_single_line = len(line_indexes) <= 1
        return runs, is_single_line

    def _estimate_single_line_height(self, elem: TextIR):
        lines = elem.lines
        line_heights = []
        if isinstance(lines, list):
            for line in lines:
                line_bbox = line.get("bbox") if isinstance(line, dict) else None
                if (
                    isinstance(line_bbox, (list, tuple))
                    and len(line_bbox) == 4
                    and line_bbox[3] > line_bbox[1]
                ):
                    line_heights.append(float(line_bbox[3] - line_bbox[1]))

        if line_heights:
            return float(np.median(line_heights))

        bbox = elem.bbox
        if isinstance(bbox, (list, tuple)) and len(bbox) == 4 and bbox[3] > bbox[1]:
            return float(bbox[3] - bbox[1])

        return 1.0

    def _compute_text_cleanup_margin_px(self, context, elem: TextIR):
        json_h = context.coords.get("json_h", 1)
        img_h = context.coords.get("img_h", 1)
        scale_y = (img_h / json_h) if json_h else 1.0

        single_line_height = self._estimate_single_line_height(elem)
        single_line_height_px = max(1.0, single_line_height * scale_y)
        return max(TEXT_CLEANUP_MIN_MARGIN_PX, int(round(single_line_height_px * self.text_cleanup_margin_ratio)))

    def _process_text(self, context, elem: TextIR):
        bbox = elem.bbox
        if not bbox:
            return

        ir_runs = elem.text_runs if isinstance(elem.text_runs, list) else []
        if not ir_runs and elem.text:
            ir_runs = [
                TextRunIR(
                    text=str(elem.text),
                    bbox=list(bbox),
                    line_index=0,
                    style=dict(elem.style or {}),
                )
            ]

        text_runs, is_single_line = self._build_text_runs_from_ir_runs(context, elem, ir_runs)
        if text_runs:
            context.add_processed_element("text", {"bbox": bbox, "text_runs": text_runs, "is_single_line": is_single_line})

    def _render_text_from_data(self, slide, text_data):
        """Renders a text element from processed data onto a slide."""
        bbox = text_data['bbox']
        text_runs = text_data.get('text_runs', [])
        is_single_line = text_data.get('is_single_line', False)

        # If it's a single line, widen the textbox to prevent wrapping due to font differences.
        if is_single_line:
            x1, y1, x2, y2 = bbox
            width = x2 - x1
            new_x2 = x1 + width * 1.2
            render_bbox = [x1, y1, new_x2, y2]
        else:
            render_bbox = bbox

        txBox = self._create_textbox(slide, render_bbox, self.coords_for_render)
        tf = txBox.text_frame
        tf.clear()
        tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = Pt(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        for run_data in text_runs:
            run = p.add_run()
            run.text = run_data['text']
            font = run.font
            font_info = run_data.get('font', {})
            font.name = font_info.get('name', "Microsoft YaHei")
            if 'size' in font_info: font.size = font_info['size']
            if 'color' in font_info: font.color.rgb = font_info['color']
            if 'bold' in font_info: font.bold = font_info['bold']

    def _add_picture_from_bbox(self, slide, bbox, page_image, coords, image_crop=None):
        if not bbox:
            return
        x1, y1, x2, y2 = bbox
        left = Pt(x1 * coords['scale_x'])
        top = Pt(y1 * coords['scale_y'])
        w = Pt((x2 - x1) * coords['scale_x'])
        h = Pt((y2 - y1) * coords['scale_y'])

        if image_crop is not None:
            crop = image_crop.copy()
        else:
            px_box = [
                int(x1 * coords['img_w'] / coords['json_w']),
                int(y1 * coords['img_h'] / coords['json_h']),
                int(x2 * coords['img_w'] / coords['json_w']),
                int(y2 * coords['img_h'] / coords['json_h']),
            ]
            crop = page_image[px_box[1]:px_box[3], px_box[0]:px_box[2]].copy()

        if crop.size > 0:
            self.temp_dir.mkdir(parents=True, exist_ok=True)
            fd, path = tempfile.mkstemp(prefix="temp_crop_img_", suffix=".png", dir=str(self.temp_dir))
            os.close(fd)
            try:
                cv2.imwrite(path, cv2.cvtColor(crop, cv2.COLOR_RGB2BGR))
                slide.shapes.add_picture(path, left, top, w, h)
            finally:
                if os.path.exists(path):
                    os.remove(path)

    def _prepare_image_crop(self, context, elem: ImageIR):
        bbox = elem.bbox
        if not bbox:
            return None
        px_box = self._to_px_bbox(bbox, context.coords)
        return context.original_image[px_box[1]:px_box[3], px_box[0]:px_box[2]].copy()

    def _cleanup_text_on_image_crop(self, context, image_elem: ImageIR, text_elem: TextIR, margin_px: int):
        if image_elem.crop_pixels is None:
            return image_elem

        txt_box = text_elem.bbox
        if not txt_box:
            return image_elem

        expanded = [
            txt_box[0] - margin_px,
            txt_box[1] - margin_px,
            txt_box[2] + margin_px,
            txt_box[3] + margin_px,
        ]
        inter = self._get_bbox_intersection(image_elem.bbox, expanded)
        if not inter:
            return image_elem

        px_img_box = self._to_px_bbox(image_elem.bbox, context.coords)
        px_inter = self._to_px_bbox(inter, context.coords)
        local_inter = [
            px_inter[0] - px_img_box[0],
            px_inter[1] - px_img_box[1],
            px_inter[2] - px_img_box[0],
            px_inter[3] - px_img_box[1],
        ]

        crop = image_elem.crop_pixels.copy()
        fill_bbox_with_bg(crop, local_inter)
        return ImageIR(
            type=image_elem.type,
            bbox=list(image_elem.bbox),
            source=image_elem.source,
            order=list(image_elem.order),
            style=dict(image_elem.style),
            is_discarded=image_elem.is_discarded,
            is_watermark=image_elem.is_watermark,
            group_id=image_elem.group_id,
            text_elements=list(image_elem.text_elements),
            crop_pixels=crop,
        )

    def _cleanup_text_and_images(self, context, text_elements: list[TextIR], image_elements: list[ImageIR]):
        text_margins: list[tuple[TextIR, int]] = []
        for text_elem in text_elements:
            margin_px = self._compute_text_cleanup_margin_px(context, text_elem)
            context.add_element_bbox_for_cleanup(text_elem.bbox, margin_px=margin_px)
            text_margins.append((text_elem, margin_px))

        cleaned_images: list[ImageIR] = []
        for image_elem in image_elements:
            current = image_elem
            for text_elem, margin_px in text_margins:
                current = self._cleanup_text_on_image_crop(context, current, text_elem, margin_px)
            cleaned_images.append(current)

        return cleaned_images

    def _process_image(self, context, elem: ImageIR):
        context.add_processed_element(
            'image',
            {
                'bbox': elem.bbox,
                'image_crop': elem.crop_pixels,
                'ir_element': elem,
            },
        )

    def _process_element(self, context, elem: ElementIR):
        if isinstance(elem, TextIR):
            self._process_text(context, elem)
        elif isinstance(elem, ImageIR):
            self._process_image(context, elem)

    def process_page(self, slide, elements, page_image, page_size=None, page_index=0, debug_images=False, context=None):
        self.debug_images = debug_images
        elements = validate_ir_elements(elements)
        elements = materialize_text_runs_for_elements(elements)
        elements = validate_ir_elements(elements, require_text_runs_consistency=True)

        img_h, img_w = page_image.shape[:2]
        json_w, json_h = page_size if page_size and all(page_size) else (img_w * 72 / 300, img_h * 72 / 300)
        w_pts, h_pts = self.cap_size(json_w, json_h)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)
        coords = {'scale_x': w_pts / json_w, 'scale_y': h_pts / json_h, 'img_w': img_w, 'img_h': img_h,
                  'json_w': json_w, 'json_h': json_h}
        self.coords_for_render = coords

        if context is None:
            context = PageContext(
                page_index,
                page_image,
                coords,
                slide,
                temp_dir=self.temp_dir,
                debug_dir=self.debug_dir,
            )
        else:
            context.slide = slide
            context.coords = coords
            context.page_index = int(page_index)

        text_elements = [e for e in elements if isinstance(e, TextIR)]
        image_elements = [e for e in elements if isinstance(e, ImageIR)]

        cleanup_text_elements = [
            e for e in text_elements if (not e.is_watermark) or self.remove_watermark
        ]
        cleanup_image_elements = [
            e for e in image_elements if (not e.is_watermark) or self.remove_watermark
        ]

        render_image_elements = [
            e for e in image_elements if not (e.is_watermark and self.remove_watermark)
        ]

        render_images_with_crop = [
            ImageIR(
                type=elem.type,
                bbox=list(elem.bbox),
                source=elem.source,
                order=list(elem.order),
                style=dict(elem.style),
                is_discarded=elem.is_discarded,
                is_watermark=elem.is_watermark,
                group_id=elem.group_id,
                text_elements=list(elem.text_elements),
                crop_pixels=self._prepare_image_crop(context, elem),
            )
            for elem in render_image_elements
        ]

        cleaned_render_images = self._cleanup_text_and_images(
            context,
            cleanup_text_elements,
            render_images_with_crop,
        )

        elements = self._normalize_page_text_font_sizes(elements)
        text_elements = [e for e in elements if isinstance(e, TextIR)]
        render_text_elements = [
            e for e in text_elements if not (e.is_watermark and self.remove_watermark)
        ]

        for image_elem in cleanup_image_elements:
            context.add_element_bbox_for_cleanup(image_elem.bbox)

        render_text_ids = {id(elem) for elem in render_text_elements}
        render_image_by_original_id = {
            id(orig): cleaned_render_images[idx]
            for idx, orig in enumerate(render_image_elements)
        }

        for elem in elements:
            if isinstance(elem, TextIR) and id(elem) in render_text_ids:
                self._process_element(context, elem)
            elif isinstance(elem, ImageIR):
                mapped = render_image_by_original_id.get(id(elem))
                if mapped is not None:
                    self._process_element(context, mapped)

        context.render_to_slide(self)

        if self.debug_images:
            context.generate_debug_images(self)

    def save(self):
        self.prs.save(self.output_path)


def _parse_pdf_page_range(page_range: str, total_pages: int) -> list[int]:
    if total_pages <= 0:
        raise ValueError("PDF has no pages to process")

    expr = str(page_range or "").strip()
    if not expr:
        raise ValueError("Page range is empty. Example: 1,3,5-8")

    selected: set[int] = set()

    for raw_token in expr.split(","):
        token = raw_token.strip()
        if not token:
            raise ValueError(f"Invalid page range token: '{raw_token}'. Example: 1,3,5-8")

        if "-" in token:
            parts = token.split("-")
            if len(parts) != 2 or not parts[0].strip().isdigit() or not parts[1].strip().isdigit():
                raise ValueError(f"Invalid page range token: '{token}'. Example: 1,3,5-8")

            start = int(parts[0].strip())
            end = int(parts[1].strip())
            if start < 1 or end < 1:
                raise ValueError(f"Page number must be >= 1, got '{token}'")
            if end < start:
                raise ValueError(f"Invalid page range '{token}': end < start")
            if end > total_pages:
                raise ValueError(f"Page number out of range in '{token}': total pages = {total_pages}")

            for page_no in range(start, end + 1):
                selected.add(page_no - 1)
        else:
            if not token.isdigit():
                raise ValueError(f"Invalid page number: '{token}'. Example: 1,3,5-8")
            page_no = int(token)
            if page_no < 1:
                raise ValueError(f"Page number must be >= 1, got '{token}'")
            if page_no > total_pages:
                raise ValueError(f"Page number out of range: '{token}', total pages = {total_pages}")
            selected.add(page_no - 1)

    return sorted(selected)


def convert_mineru_to_ppt(
    json_path,
    input_path,
    output_ppt_path,
    remove_watermark=True,
    debug_images=False,
    ocr_engine=None,
    ocr_device_policy="auto",
    ocr_model_root=None,
    ocr_model_variant="auto",
    ocr_offline_only=False,
    ocr_det_db_thresh=None,
    ocr_det_db_box_thresh=None,
    ocr_det_db_unclip_ratio=None,
    page_range=None,
    text_cleanup_margin_ratio=None,
    ocr_font_distance_threshold=None,
):
    from .utils import pdf_to_images
    DPI = 300

    output_dir = Path(output_ppt_path).resolve().parent
    debug_dir = output_dir / "debug"
    temp_dir = debug_dir if debug_images else Path(tempfile.gettempdir())

    if debug_images:
        if debug_dir.exists():
            shutil.rmtree(debug_dir)
        debug_dir.mkdir(parents=True, exist_ok=True)

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    is_pdf_input = input_path.lower().endswith('.pdf')
    if is_pdf_input:
        images = pdf_to_images(input_path, dpi=DPI)
    else:
        try:
            img = Image.open(input_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            images = [np.array(img)]
        except Exception as e:
            raise IOError(f"Failed to load image file: {input_path} - {e}")

    if ocr_engine is None:
        ocr_engine = PaddleOCREngine(
            device_policy=ocr_device_policy,
            use_angle_cls=False,
            model_root=ocr_model_root,
            offline_only=ocr_offline_only,
            det_db_thresh=ocr_det_db_thresh,
            det_db_box_thresh=ocr_det_db_box_thresh,
            det_db_unclip_ratio=ocr_det_db_unclip_ratio,
            refine_font_distance_threshold=ocr_font_distance_threshold,
            model_variant=ocr_model_variant,
        )
    elif ocr_font_distance_threshold is not None:
        ocr_engine.refine_font_distance_threshold = ocr_font_distance_threshold

    mineru_adapter = MinerUAdapter()
    ocr_adapter = OCRAdapter(ocr_engine)

    gen = PPTGenerator(
        output_ppt_path,
        remove_watermark=remove_watermark,
        ocr_engine=ocr_engine,
        ocr_device_policy=ocr_device_policy,
        ocr_model_root=ocr_model_root,
        ocr_model_variant=ocr_model_variant,
        ocr_offline_only=ocr_offline_only,
        text_cleanup_margin_ratio=text_cleanup_margin_ratio,
        temp_dir=temp_dir,
        debug_dir=(debug_dir if debug_images else None),
    )
    pages = data if isinstance(data, list) else next(
        (data[k] for k in ["pdf_info", "pages"] if k in data and isinstance(data[k], list)), [data]
    )
    print(f"[CLEANUP] Found {len(pages)} pages.")

    selected_page_indices = list(range(len(pages)))
    if is_pdf_input and page_range is not None and str(page_range).strip() != "":
        selected_page_indices = _parse_pdf_page_range(str(page_range), len(pages))

    first_processed = True
    total_selected = len(selected_page_indices)

    for selected_pos, page_index in enumerate(selected_page_indices):
        if page_index >= len(images):
            break

        page_data = pages[page_index]
        print(f"Processing page {page_index + 1}/{len(pages)} (selected {selected_pos + 1}/{total_selected})...")
        page_img = images[page_index].copy()

        if first_processed:
            gen.set_slide_size(page_img.shape[1], page_img.shape[0], dpi=DPI)
            first_processed = False

        slide = gen.add_slide()

        page_size = page_data.get("page_size") or (
            page_data.get("page_info", {}).get("width"),
            page_data.get("page_info", {}).get("height"),
        )
        json_w, json_h = page_size if page_size and all(page_size) else (page_img.shape[1], page_img.shape[0])
        coords = {
            'scale_x': 1.0,
            'scale_y': 1.0,
            'img_w': page_img.shape[1],
            'img_h': page_img.shape[0],
            'json_w': json_w,
            'json_h': json_h,
        }
        page_context = PageContext(
            page_index,
            page_img,
            coords,
            slide,
            temp_dir=temp_dir,
            debug_dir=(debug_dir if debug_images else None),
        )

        mineru_page = MinerUPageData.from_dict(page_data)
        mineru_elements = validate_ir_elements(mineru_adapter.extract_page_elements(mineru_page, include_text_runs=False))
        page_context.register_stage_page_ir(
            "mineru_original",
            build_page_ir(page_index=page_index, page_size=(float(json_w), float(json_h)), elements=mineru_elements),
        )

        try:
            ocr_elements = validate_ir_elements(
                ocr_adapter.extract_page_elements(page_img, json_w, json_h, page_context=page_context)
            )
        except Exception as e:
            raise RuntimeError(f"[OCR] Page {page_index + 1}: OCR extraction failed: {e}") from e

        merged_elements, merge_stats = merge_ir_elements(mineru_elements, ocr_elements, gen._has_bbox_overlap)
        page_context.register_stage_page_ir(
            "merged_final",
            build_page_ir(page_index=page_index, page_size=(float(json_w), float(json_h)), elements=merged_elements),
        )
        print(
            f"[OCR] Page {page_index + 1}: "
            f"candidates={merge_stats['overlay_candidates']}, "
            f"group_replaced={merge_stats['group_replaced']}, "
            f"overlap_replaced={merge_stats['overlap_replaced']}, "
            f"added={merge_stats['overlay_added']}"
        )

        gen.process_page(
            slide,
            merged_elements,
            page_img,
            page_size=page_size,
            page_index=page_index,
            debug_images=debug_images,
            context=page_context,
        )
    gen.save()
    print(f"Saved to {output_ppt_path}")
