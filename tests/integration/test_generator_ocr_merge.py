import json
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import numpy as np
from PIL import Image
from pptx import Presentation

from converter.generator import PPTGenerator, convert_mineru_to_ppt, PageContext
from converter.ir import TextIR, TextRunIR
from converter.ir_merge import merge_ir_elements
from converter.ocr_merge import PaddleOCREngine


class FakeOCREngine:
    def __init__(self, should_fail=False):
        self.should_fail = should_fail

    def extract_text_elements(self, page_image, json_w, json_h, return_stage_elements=False):
        if self.should_fail:
            raise RuntimeError("ocr init failed")
        elements = [
            {
                "type": "text",
                "bbox": [60, 60, 100, 100],
                "index": 1,
                "lines": [
                    {
                        "bbox": [60, 60, 100, 100],
                        "spans": [{"bbox": [60, 60, 100, 100], "content": "ocr-text", "type": "text"}],
                    }
                ],
                "is_discarded": False,
            }
        ]
        if return_stage_elements:
            return {
                "before_refined_elements": elements,
                "after_refined_elements": elements,
            }
        return elements


class TestGeneratorOCRMerge(unittest.TestCase):
    def test_process_page_with_forced_ocr_does_not_raise(self):
        generator = PPTGenerator("out.pptx", ocr_engine=FakeOCREngine())
        slide = generator.add_slide()
        page_image = np.zeros((200, 200, 3), dtype=np.uint8)

        elements = [
            {
                "type": "text",
                "bbox": [10, 10, 40, 40],
                "text": "mineru",
                "is_discarded": False,
            }
        ]

        generator.process_page(slide, elements, page_image, page_size=(200, 200), page_index=0, debug_images=False)

    def test_convert_function_ocr_failure_raises(self):
        fake_json_data = [{"para_blocks": [], "images": [], "tables": [], "discarded_blocks": []}]

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.generator.Image.open") as mocked_image_open,
            mock.patch("converter.generator.np.array", return_value=np.zeros((600, 600, 3), dtype=np.uint8)),
            mock.patch("converter.generator.PPTGenerator") as mocked_generator_cls,
        ):
            rgb_image = mock.Mock()
            rgb_image.mode = "RGB"
            rgb_image.convert.return_value = rgb_image
            mocked_image_open.return_value = rgb_image

            generator_instance = mocked_generator_cls.return_value
            generator_instance.add_slide.return_value = object()

            with self.assertRaises(RuntimeError):
                convert_mineru_to_ppt(
                    "in.json",
                    "in.png",
                    "out.pptx",
                    ocr_engine=FakeOCREngine(should_fail=True),
                )

    def test_convert_applies_pdf_page_range(self):
        fake_json_data = [
            {
                "para_blocks": [{"type": "text", "bbox": [10, 10, 90, 40], "text": "p1"}],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            },
            {
                "para_blocks": [{"type": "text", "bbox": [10, 10, 90, 40], "text": "p2"}],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            },
            {
                "para_blocks": [{"type": "text", "bbox": [10, 10, 90, 40], "text": "p3"}],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            },
        ]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        fake_images = [np.zeros((600, 600, 3), dtype=np.uint8) for _ in range(3)]

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.utils.pdf_to_images", return_value=fake_images),
            mock.patch.object(PPTGenerator, "save", return_value=None),
        ):
            gen = convert_mineru_to_ppt(
                "in.json",
                "in.pdf",
                "out.pptx",
                ocr_engine=_Engine(),
                page_range="1,3",
            )

            self.assertIsNone(gen)

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.utils.pdf_to_images", return_value=fake_images),
            mock.patch.object(PPTGenerator, "save", return_value=None),
        ):
            generator = PPTGenerator("out.pptx", ocr_engine=_Engine())
            with mock.patch("converter.generator.PPTGenerator", return_value=generator):
                convert_mineru_to_ppt(
                    "in.json",
                    "in.pdf",
                    "out.pptx",
                    ocr_engine=_Engine(),
                    page_range="1,3",
                )
                self.assertEqual(len(generator.prs.slides), 2)

    def test_convert_rejects_invalid_pdf_page_range(self):
        fake_json_data = [
            {
                "para_blocks": [{"type": "text", "bbox": [10, 10, 90, 40], "text": "p1"}],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            }
        ]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.utils.pdf_to_images", return_value=[np.zeros((600, 600, 3), dtype=np.uint8)]),
            mock.patch.object(PPTGenerator, "save", return_value=None),
        ):
            with self.assertRaises(ValueError):
                convert_mineru_to_ppt(
                    "in.json",
                    "in.pdf",
                    "out.pptx",
                    ocr_engine=_Engine(),
                    page_range="2",
                )

    def test_convert_ignores_page_range_for_non_pdf(self):
        fake_json_data = [
            {
                "para_blocks": [{"type": "text", "bbox": [10, 10, 90, 40], "text": "img"}],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            }
        ]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        with (
            tempfile.TemporaryDirectory() as td,
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch.object(PPTGenerator, "save", return_value=None),
        ):
            input_png = Path(td) / "in.png"
            output_ppt = Path(td) / "out.pptx"
            Image.fromarray(np.zeros((600, 600, 3), dtype=np.uint8)).save(input_png)

            convert_mineru_to_ppt(
                "in.json",
                str(input_png),
                str(output_ppt),
                ocr_engine=_Engine(),
                page_range="999",
            )

    def test_convert_constructs_ocr_engine_with_forwarded_config(self):
        fake_json_data = [{"para_blocks": [], "images": [], "tables": [], "discarded_blocks": []}]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.generator.Image.open") as mocked_image_open,
            mock.patch("converter.generator.np.array", return_value=np.zeros((600, 600, 3), dtype=np.uint8)),
            mock.patch("converter.generator.PaddleOCREngine", return_value=_Engine()) as mocked_engine,
            mock.patch("converter.generator.PPTGenerator") as mocked_generator_cls,
        ):
            rgb_image = mock.Mock()
            rgb_image.mode = "RGB"
            rgb_image.convert.return_value = rgb_image
            mocked_image_open.return_value = rgb_image

            generator_instance = mocked_generator_cls.return_value
            generator_instance.add_slide.return_value = object()

            convert_mineru_to_ppt(
                "in.json",
                "in.png",
                "out.pptx",
                ocr_device_policy="cpu",
                ocr_model_root="models/paddleocr",
                ocr_model_variant="server",
                ocr_offline_only=False,
            )

        mocked_engine.assert_called_once_with(
            device_policy="cpu",
            use_angle_cls=False,
            model_root="models/paddleocr",
            offline_only=False,
            det_db_thresh=None,
            det_db_box_thresh=None,
            det_db_unclip_ratio=None,
            refine_font_distance_threshold=None,
            model_variant="server",
        )

    def test_convert_forwards_tuned_ocr_detection_options_to_engine(self):
        fake_json_data = [{"para_blocks": [], "images": [], "tables": [], "discarded_blocks": []}]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.generator.Image.open") as mocked_image_open,
            mock.patch("converter.generator.np.array", return_value=np.zeros((600, 600, 3), dtype=np.uint8)),
            mock.patch("converter.generator.PaddleOCREngine", return_value=_Engine()) as mocked_engine,
            mock.patch("converter.generator.PPTGenerator") as mocked_generator_cls,
        ):
            rgb_image = mock.Mock()
            rgb_image.mode = "RGB"
            rgb_image.convert.return_value = rgb_image
            mocked_image_open.return_value = rgb_image

            generator_instance = mocked_generator_cls.return_value
            generator_instance.add_slide.return_value = object()

            convert_mineru_to_ppt(
                "in.json",
                "in.png",
                "out.pptx",
                ocr_det_db_thresh=0.35,
                ocr_det_db_box_thresh=0.8,
                ocr_det_db_unclip_ratio=1.1,
            )

        mocked_engine.assert_called_once_with(
            device_policy="auto",
            use_angle_cls=False,
            model_root=None,
            offline_only=False,
            det_db_thresh=0.35,
            det_db_box_thresh=0.8,
            det_db_unclip_ratio=1.1,
            refine_font_distance_threshold=None,
            model_variant="auto",
        )

    def test_convert_function_forwards_ocr_config_to_generator(self):
        fake_json_data = [{"para_blocks": [], "images": [], "tables": [], "discarded_blocks": []}]

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return []

        with (
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.generator.Image.open") as mocked_image_open,
            mock.patch("converter.generator.PPTGenerator") as mocked_generator_cls,
            mock.patch("converter.generator.np.array", return_value=np.zeros((600, 600, 3), dtype=np.uint8)),
            mock.patch("converter.generator.PaddleOCREngine", return_value=_Engine()),
        ):
            rgb_image = mock.Mock()
            rgb_image.mode = "RGB"
            rgb_image.convert.return_value = rgb_image
            mocked_image_open.return_value = rgb_image

            generator_instance = mocked_generator_cls.return_value
            generator_instance.add_slide.return_value = object()

            convert_mineru_to_ppt(
                "in.json",
                "in.png",
                "out.pptx",
                ocr_device_policy="gpu",
                ocr_model_root="x/models/paddleocr",
                ocr_model_variant="server",
                ocr_offline_only=False,
            )

        self.assertEqual(mocked_generator_cls.call_count, 1)
        kwargs = mocked_generator_cls.call_args.kwargs
        self.assertEqual(kwargs["remove_watermark"], True)
        self.assertEqual(kwargs["ocr_device_policy"], "gpu")
        self.assertEqual(kwargs["ocr_model_root"], "x/models/paddleocr")
        self.assertEqual(kwargs["ocr_model_variant"], "server")
        self.assertEqual(kwargs["ocr_offline_only"], False)
        self.assertIsNotNone(kwargs["ocr_engine"])


    def test_convert_outputs_stage_debug_images_and_original_png(self):
        fake_json_data = [
            {
                "para_blocks": [
                    {
                        "type": "text",
                        "bbox": [10, 10, 90, 40],
                        "text": "mineru-text",
                    }
                ],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            }
        ]

        fake_ocr_output = [
            {
                "type": "text",
                "bbox": [10, 10, 90, 40],
                "index": 1,
                "lines": [
                    {
                        "bbox": [10, 10, 90, 40],
                        "spans": [
                            {
                                "bbox": [10, 10, 90, 40],
                                "content": "ocr-text",
                                "type": "text",
                            }
                        ],
                    }
                ],
                "is_discarded": False,
            }
        ]

        with (
            tempfile.TemporaryDirectory() as td,
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
            mock.patch("converter.generator.os.path.exists", return_value=False),
            mock.patch("converter.generator.os.makedirs", return_value=None),
            mock.patch("converter.generator.shutil.rmtree", return_value=None),
            mock.patch("converter.generator.cv2.imwrite", return_value=True) as mocked_imwrite,
            mock.patch("converter.generator.PageContext.render_to_slide", return_value=None),
        ):
            output_ppt = Path(td) / "out.pptx"
            input_png = Path(td) / "in.png"
            Image.fromarray(np.zeros((600, 600, 3), dtype=np.uint8)).save(input_png)

            class _Engine:
                def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                    if return_stage_elements:
                        return {
                            "before_refined_elements": fake_ocr_output,
                            "after_refined_elements": fake_ocr_output,
                        }
                    return fake_ocr_output

            convert_mineru_to_ppt(
                "in.json",
                str(input_png),
                str(output_ppt),
                ocr_engine=_Engine(),
                debug_images=True,
            )

            written_paths = [
                str(call.args[0]).replace("\\", "/")
                for call in mocked_imwrite.call_args_list
                if call.args
            ]
            self.assertTrue(any(path.endswith("/debug/page_0_original.png") for path in written_paths))
            self.assertTrue(any(path.endswith("/debug/page_0_mineru_original_boxes.png") for path in written_paths))
            self.assertTrue(any(path.endswith("/debug/page_0_ocr_before_refined_elements.png") for path in written_paths))
            self.assertTrue(any(path.endswith("/debug/page_0_ocr_after_refined_elements.png") for path in written_paths))
            self.assertTrue(any(path.endswith("/debug/page_0_merged_final_boxes.png") for path in written_paths))

    def test_convert_inherits_mineru_bold_when_ocr_replaces_text(self):
        fake_json_data = [
            {
                "para_blocks": [
                    {
                        "type": "title",
                        "bbox": [10, 10, 90, 40],
                        "text": "mineru-title",
                    }
                ],
                "images": [],
                "tables": [],
                "discarded_blocks": [],
                "page_size": [600, 600],
            }
        ]

        fake_ocr_output = [
            {
                "type": "text",
                "bbox": [10, 10, 90, 40],
                "index": 1,
                "lines": [
                    {
                        "bbox": [10, 10, 90, 40],
                        "spans": [
                            {
                                "bbox": [10, 10, 90, 40],
                                "content": "ocr-title",
                                "type": "text",
                            }
                        ],
                    }
                ],
                "is_discarded": False,
            }
        ]

        with (
            tempfile.TemporaryDirectory() as td,
            mock.patch("converter.generator.open", mock.mock_open(read_data="[]")),
            mock.patch("converter.generator.json.load", return_value=fake_json_data),
        ):
            output_ppt = Path(td) / "out.pptx"
            input_png = Path(td) / "in.png"
            Image.fromarray(np.zeros((600, 600, 3), dtype=np.uint8)).save(input_png)

            class _Engine:
                def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                    if return_stage_elements:
                        return {
                            "before_refined_elements": fake_ocr_output,
                            "after_refined_elements": fake_ocr_output,
                        }
                    return fake_ocr_output

            convert_mineru_to_ppt(
                "in.json",
                str(input_png),
                str(output_ppt),
                ocr_engine=_Engine(),
            )

            prs = Presentation(str(output_ppt))
            slide = prs.slides[0]

            text_runs = []
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                tf = shape.text_frame
                if tf is None:
                    continue
                for para in tf.paragraphs:
                    text_runs.extend(para.runs)

            self.assertTrue(any(run.text == "ocr-title" for run in text_runs))
            self.assertTrue(
                any(run.text == "ocr-title" and run.font.bold for run in text_runs),
                "Expected OCR-replaced text to inherit MinerU bold in rendered PPT",
            )

    def test_process_page_runs_font_normalization_before_render(self):
        generator = PPTGenerator("out.pptx", ocr_engine=FakeOCREngine())
        slide = generator.add_slide()
        page_image = np.zeros((200, 200, 3), dtype=np.uint8)
        elements = [
            {
                "type": "text",
                "bbox": [10, 10, 50, 40],
                "text": "sample",
                "is_discarded": False,
            }
        ]

        events = []

        def _normalize_hook(page_elements):
            events.append("normalize")
            return page_elements

        def _render_hook(_generator):
            events.append("render")
            return None

        with (
            mock.patch.object(generator, "_normalize_page_text_font_sizes", side_effect=_normalize_hook) as mocked_normalize,
            mock.patch("converter.generator.PageContext.render_to_slide", side_effect=_render_hook),
        ):
            generator.process_page(slide, elements, page_image, page_size=(200, 200), page_index=0, debug_images=False)

        self.assertEqual(mocked_normalize.call_count, 1)
        self.assertEqual(events, ["normalize", "render"])

    def test_process_page_preserves_image_bbox_after_font_normalization_stage(self):
        generator = PPTGenerator("out.pptx", ocr_engine=FakeOCREngine())
        slide = generator.add_slide()
        page_image = np.zeros((200, 200, 3), dtype=np.uint8)

        elements = [
            {
                "type": "text",
                "bbox": [10, 10, 50, 40],
                "text": "caption",
                "is_discarded": False,
            },
            {
                "type": "image",
                "bbox": [120, 100, 180, 160],
                "is_discarded": False,
            },
        ]

        coords = {
            "scale_x": 1.0,
            "scale_y": 1.0,
            "img_w": 200,
            "img_h": 200,
            "json_w": 200,
            "json_h": 200,
        }
        context = PageContext(page_index=0, page_image=page_image, coords=coords, slide=slide)

        captured_elements = []

        original_render = PageContext.render_to_slide

        def _render_hook(patched_context, patched_generator):
            captured_elements.extend(patched_context.elements)
            return original_render(patched_context, patched_generator)

        with mock.patch("converter.generator.PageContext.render_to_slide", autospec=True, side_effect=_render_hook):
            generator.process_page(
                slide,
                elements,
                page_image,
                page_size=(200, 200),
                page_index=0,
                debug_images=False,
                context=context,
            )

        image_entries = [entry for entry in captured_elements if entry.get("type") == "image"]
        self.assertEqual(len(image_entries), 1)
        self.assertEqual(image_entries[0]["data"]["bbox"], [120.0, 100.0, 180.0, 160.0])

    def test_case2_ir_merge_outputs_text_elements(self):
        repo_root = Path(__file__).resolve().parents[2]
        input_png = repo_root / "demo" / "case2" / "PixPin_2026-03-05_22-01-24.png"
        input_json = repo_root / "demo" / "case2" / "MinerU_PixPin_2026-03-05_22-01-24__20260305140239.json"

        self.assertTrue(input_png.exists(), f"Missing demo image: {input_png}")
        self.assertTrue(input_json.exists(), f"Missing demo json: {input_json}")

        data = json.loads(input_json.read_text(encoding="utf-8"))
        page_data = data["pdf_info"][0]

        elements = []
        for item in page_data.get("para_blocks", []):
            if item:
                elements.append(item)
        for item in page_data.get("discarded_blocks", []):
            if item:
                item["is_discarded"] = True
                elements.append(item)

        page_image = np.array(Image.open(input_png).convert("RGB"))

        class _Engine:
            def extract_text_elements(self, *_args, return_stage_elements=False, **_kwargs):
                if return_stage_elements:
                    return {
                        "before_refined_elements": [],
                        "after_refined_elements": [],
                    }
                return [
                    {
                        "type": "text",
                        "bbox": [10.0, 10.0, 90.0, 40.0],
                        "index": 1,
                        "lines": [
                            {
                                "bbox": [10.0, 10.0, 90.0, 40.0],
                                "spans": [
                                    {
                                        "bbox": [10.0, 10.0, 90.0, 40.0],
                                        "content": "构建流水线",
                                        "type": "text",
                                    }
                                ],
                            }
                        ],
                        "is_discarded": False,
                    },
                    {
                        "type": "text",
                        "bbox": [12.0, 50.0, 210.0, 80.0],
                        "index": 2,
                        "lines": [
                            {
                                "bbox": [12.0, 50.0, 210.0, 80.0],
                                "spans": [
                                    {
                                        "bbox": [12.0, 50.0, 210.0, 80.0],
                                        "content": "形成统一规范的数据开发标准",
                                        "type": "text",
                                    }
                                ],
                            }
                        ],
                        "is_discarded": False,
                    },
                ]

        ocr_engine = _Engine()
        ocr_elements = ocr_engine.extract_text_elements(
            page_image,
            page_image.shape[1],
            page_image.shape[0],
        )

        def _has_overlap(a, b):
            x1 = max(a[0], b[0])
            y1 = max(a[1], b[1])
            x2 = min(a[2], b[2])
            y2 = min(a[3], b[3])
            return x1 < x2 and y1 < y2

        # simulate adapter output contract for this integration case
        mineru_ir: list[TextIR] = []
        for elem in elements:
            if not elem.get("bbox"):
                continue
            lines = elem.get("lines", [])
            text_from_lines = "\n".join(
                "".join(span.get("content", "") for span in line.get("spans", []))
                for line in lines
                if line.get("spans")
            )
            text_value = text_from_lines or elem.get("text", "")
            if not text_value and not lines:
                continue
            mineru_ir.append(
                TextIR(
                    type="text",
                    bbox=[float(v) for v in elem.get("bbox")],
                    lines=lines,
                    text=str(text_value),
                    group_id=elem.get("group_id"),
                    order=[float(elem.get("index", 0)), 0.0],
                    style={"bold": False, "font_size": None, "align": "left"},
                    is_discarded=bool(elem.get("is_discarded", False)),
                    source="mineru",
                    text_runs=None,
                )
            )

        ocr_ir: list[TextIR] = []
        for elem in ocr_elements:
            if not elem.get("bbox"):
                continue
            lines = elem.get("lines", [])
            text = "\n".join(
                "".join(span.get("content", "") for span in line.get("spans", []))
                for line in lines
                if line.get("spans")
            )
            runs = [
                TextRunIR(
                    text=str(span.get("content", "")),
                    bbox=[float(v) for v in (span.get("bbox") or line.get("bbox") or elem.get("bbox"))],
                    line_index=line_index,
                    style={},
                )
                for line_index, line in enumerate(lines)
                for span in line.get("spans", [])
                if span.get("content", "")
            ]
            ocr_ir.append(
                TextIR(
                    type="text",
                    bbox=[float(v) for v in elem.get("bbox")],
                    lines=lines,
                    text=text,
                    text_runs=runs,
                    group_id=elem.get("group_id"),
                    order=[float(elem.get("index", 0)), 0.0],
                    style={"bold": False, "font_size": None, "align": "left"},
                    is_discarded=False,
                    source="ocr",
                )
            )

        merged, _stats = merge_ir_elements(mineru_ir, ocr_ir, _has_overlap)

        merged_texts = [
            str(elem.text)
            for elem in merged
            if isinstance(elem, TextIR)
        ]

        self.assertTrue(any("构建" in text and "流水线" in text for text in merged_texts))
        self.assertTrue(any("形成统一规范的数据开发标准" in text for text in merged_texts))

        target_phrase_elements = [
            elem
            for elem in merged
            if isinstance(elem, TextIR)
            and "设计即开发" in str(elem.text)
            and "十倍提效" in str(elem.text)
        ]
        self.assertEqual(
            len(target_phrase_elements),
            1,
            "Expected sentence '以“设计即开发”实现开发侧十倍提效。' to stay in a single TextIR",
        )

if __name__ == "__main__":
    unittest.main()
